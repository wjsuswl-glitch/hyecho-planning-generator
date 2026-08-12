"""동적 PPTX 빌더 (v2) — 기존 완성 기획안을 열어서 도형을 덮어쓰는 대신,
매번 필요한 만큼만 슬라이드를 새로 만든다.

기존 assembler.py 방식(템플릿 기획안을 열어 shape_id별로 텍스트 주입)은
- 옛 기획안의 이미지/깨진 도형이 그대로 남는 문제
- 목적지 개수가 템플릿 슬롯 수와 안 맞으면 삭제 로직이 계속 필요한 문제
가 있어서, 이 모듈은 옛 기획안(예: 안데스_템플릿.pptx)을 "디자인 참고용"으로만
쓰고, 실제 출력은 python-pptx로 슬라이드를 새로 그린다.
"""
import json
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- 스타일 상수 (정현지 스타일 기본값) ----
SLIDE_W = Emu(6858000)   # 원본 기획안과 동일한 슬라이드 크기 (약 7.5 x 10.83 inch, 모바일 세로형)
SLIDE_H = Emu(9906000)
FONT_NAME = "맑은 고딕"
ACCENT_COLOR = RGBColor(0x1B, 0x4D, 0x6B)   # 진한 티얼/네이비 (섹션 바, 강조)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_COLOR = RGBColor(0x66, 0x66, 0x66)
MARGIN = Inches(0.4)
CONTENT_W = SLIDE_W - MARGIN * 2


def _tf_setup(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=FONT_NAME):
    tf.word_wrap = True
    lines = str(text).split("\n")
    tf.text = lines[0]
    p0 = tf.paragraphs[0]
    p0.alignment = align
    for run in p0.runs or [p0.add_run()]:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.alignment = align
        r = p.add_run() if not p.runs else p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font


def estimate_text_height(text, size_pt, width_emu, line_spacing=1.22, bold=False):
    """글자 수 기반으로 텍스트가 실제로 차지할 높이를 대략 추정한다.
    고정 간격 대신 이걸 써야 설명 길이에 따라 다음 요소와 안 겹친다.
    실제 렌더링 폭 추정은 부정확할 수 있어 여유 마진을 둔다(단, 슬라이드 장수를
    압축하기 위해 과도한 여유는 줄였다 — line_spacing 1.35→1.22, 버퍼 0.15→0.08in)."""
    if not text:
        return Inches(0.1)
    width_in = Emu(width_emu).inches
    # 한글 기준 글자 폭 대략치 (볼드면 좀 더 넓게 잡음), 안전 마진 포함
    char_w_in = (size_pt / 72) * (1.05 if bold else 0.95)
    chars_per_line = max(1, int(width_in / char_w_in))
    total_lines = 0
    for line in str(text).split("\n"):
        total_lines += max(1, -(-len(line) // chars_per_line))  # ceil
    line_h_in = (size_pt / 72) * line_spacing
    return Inches(total_lines * line_h_in + 0.08)  # 여유 마진


def add_text(slide, left, top, width, height, text, size=14, color=TEXT_COLOR,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    _tf_setup(tf, text, size, color, bold, align)
    return box


def add_section_bar(slide, top, text, height=Inches(0.45), size=15):
    """섹션 제목이 들어가는 진한 색 배경 바 (예: '천장공로 하이라이트')"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top, CONTENT_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _tf_setup(tf, text, size, WHITE, bold=True, align=PP_ALIGN.CENTER)
    return bar


def add_image_placeholder(slide, left, top, width, height, label="이미지"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.background()
    box.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _tf_setup(tf, label, 11, MUTED_COLOR, align=PP_ALIGN.CENTER)
    return box


def add_small_image_placeholder(slide, top, width, height, label="이미지"):
    """전체 폭을 다 차지하는 큰 이미지 자리 대신, 가로 폭도 훨씬 좁힌 작은 썸네일
    크기 자리표시. CONTENT_W 안에서 가운데 정렬해서 배치한다. 실제 사진은 디자이너가
    작업하므로, 여기서는 '사진이 들어갈 위치'만 작게 표시하고 나머지 공간은 텍스트가
    위로 당겨져 채우게 한다."""
    left = MARGIN + (CONTENT_W - width) / 2
    return add_image_placeholder(slide, left, top, width, height, label)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 완전 빈 레이아웃


class SlideFlow:
    """여러 섹션 함수가 슬라이드 하나를 공유해서 이어 쓸 수 있게 하는 커서.

    PPTX는 프레젠테이션 전체가 슬라이드 높이를 하나만 가질 수 있어서(슬라이드별
    높이 지정 불가), 식사 안내(1.6in)처럼 짧은 섹션도 그동안은 10.83in짜리 슬라이드를
    통째로 차지해 빈 공간이 컸다. ensure()로 "남은 공간에 들어가면 이어 붙이고,
    안 들어가면 새 슬라이드"를 섹션마다 판단해 슬라이드 장수와 여백을 줄인다."""

    def __init__(self, prs):
        self.prs = prs
        self.slide = None
        self.y = Inches(0.3)
        self.bottom_limit = SLIDE_H - Inches(0.2)

    def new_slide(self):
        self.slide = _blank_slide(self.prs)
        self.y = Inches(0.3)
        return self.slide

    def ensure(self, height, gap_before=Inches(0.35)):
        """다음 섹션(height)을 놓을 자리를 확보한다. 이어 붙일 수 있으면 그 y좌표를,
        없으면 새 슬라이드를 시작하고 그 y좌표를 반환한다."""
        if self.slide is None:
            self.new_slide()
            return self.y
        candidate_y = self.y + gap_before
        if candidate_y + height > self.bottom_limit:
            self.new_slide()
            return self.y
        self.y = candidate_y
        return self.y


# ---------------------------------------------------------------------------
# 슬라이드 빌더 함수 (정현지 스타일) — 모두 SlideFlow를 받아 가능하면 같은
# 슬라이드에 이어 그리고, 공간이 없을 때만 새 슬라이드를 시작한다.
# ---------------------------------------------------------------------------

def build_cover_slide(flow, cover, watermark_label=""):
    """표지 슬라이드. tagline/product_name/subtitle 모두 고정 높이를 확보해두고
    있었는데, 이 셋 다 AI가 채우는 가변 길이 텍스트라 길어지면 예상보다 줄바꿈이
    늘어 바로 아래 요소와 겹칠 수 있다(background_story에서 실제로 발생한 것과
    같은 문제 — 표지는 모든 상품에서 항상 렌더링되는 슬라이드라 잠재 위험이 가장
    크다). estimate_text_height로 실제 줄 수를 추정해 안전하게 확보한다."""
    slide = flow.new_slide()
    y = flow.y
    tagline = cover.get("tagline", "")
    tagline_h = estimate_text_height(tagline, 13, CONTENT_W) if tagline else Inches(0.5)
    add_text(slide, MARGIN, y, CONTENT_W, tagline_h, tagline,
              size=13, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
    y += tagline_h + Inches(0.05)

    product_name = cover.get("product_name", "")
    product_h = estimate_text_height(product_name, 26, CONTENT_W, bold=True) if product_name else Inches(0.9)
    add_text(slide, MARGIN, y, CONTENT_W, product_h, product_name,
              size=26, bold=True, align=PP_ALIGN.CENTER)
    y += product_h + Inches(0.05)

    if cover.get("region_tag"):
        region_h = estimate_text_height(cover["region_tag"], 13, CONTENT_W)
        add_text(slide, MARGIN, y, CONTENT_W, region_h, cover["region_tag"],
                  size=13, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
        y += region_h + Inches(0.05)
    y += Inches(0.1)
    add_small_image_placeholder(slide, y, Inches(2.3), Inches(0.9), "메인 이미지")
    y += Inches(1.0)
    if cover.get("subtitle"):
        subtitle_h = estimate_text_height(cover["subtitle"], 14, CONTENT_W, bold=True)
        add_text(slide, MARGIN, y, CONTENT_W, subtitle_h, cover["subtitle"],
                  size=14, bold=True, align=PP_ALIGN.CENTER)
        y += subtitle_h + Inches(0.05)
    intro_h = estimate_text_height(cover.get("intro_copy", ""), 12, CONTENT_W)
    add_text(slide, MARGIN, y, CONTENT_W, intro_h, cover.get("intro_copy", ""),
              size=12, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
    y += intro_h
    if watermark_label:
        add_text(slide, SLIDE_W - Inches(1.5), Inches(0.15), Inches(1.1), Inches(0.3),
                  watermark_label, size=11, bold=True, color=RGBColor(0xCC, 0xB0, 0x00),
                  align=PP_ALIGN.RIGHT)
    flow.y = y
    return slide


def build_destination_slides(flow, destinations, section_title=None, theme_line=None):
    """목적지 개수만큼만 슬라이드를 만든다. 고정 개수로 나누지 않고, 실제 텍스트
    길이를 추정해서 한 슬라이드에 들어갈 수 있는 만큼만 채우고 넘치면 다음 슬라이드로.
    첫 슬라이드는 이전 섹션이 남긴 여백에 이어 붙일 수 있으면 이어 붙인다.

    region_tag가 있는 방문지는 같은 그룹(성격이 비슷하거나 같은 지역)끼리 묶어서
    그룹명을 소제목으로 한 번만 보여준다 — 예전엔 방문지마다 작은 라벨을 하나씩
    달았는데, 그 라벨이 흰 글씨(WHITE)를 배경 없이 그냥 텍스트로 찍어서 화면에서
    보이지도 않았고(멕시코 문명기행 테스트로 확인), 설령 보였어도 방문지 하나하나에
    태그만 붙을 뿐 "묶어서 보여주는" 느낌은 아니었다. 사업부 자료에 이미 있는 그룹
    구조(예: "[ 사포테카 문명 | 와하카 ]")를 살려 지역/테마별로 소제목 아래 방문지를
    들여쓰기해서 묶어 보여주도록 다시 짰다 — region_tag가 없는 방문지는 예전처럼
    그룹 소제목 없이 그냥 나열된다(하위 호환)."""
    if not destinations:
        return []
    slides = []
    idx = 0
    first_slide = True
    header_h = (Inches(0.55) if section_title else Inches(0)) + \
               (Inches(0.4) if theme_line else Inches(0))
    current_group = None  # 화면에 마지막으로 그린 그룹명 — 슬라이드가 넘어가도 유지

    while idx < len(destinations):
        if first_slide:
            y = flow.ensure(header_h + Inches(0.5))  # 헤더 + 항목 하나 들어갈 여유
            slide = flow.slide
            if section_title:
                add_section_bar(slide, y, section_title)
                y += Inches(0.55)
            if theme_line:
                add_text(slide, MARGIN, y, CONTENT_W, Inches(0.35), theme_line,
                          size=14, bold=True, align=PP_ALIGN.CENTER)
                y += Inches(0.4)
            first_slide = False
        else:
            slide = flow.new_slide()
            y = flow.y
            # 그룹이 슬라이드 경계를 넘어 이어지면, 새 슬라이드 맨 위에서도 그룹명을
            # 다시 보여준다 — 안 그러면 방문지가 갑자기 그룹 없이 나온 것처럼 보인다.
            if current_group:
                label_h = estimate_text_height(current_group, 13, CONTENT_W, bold=True)
                add_text(slide, MARGIN, y, CONTENT_W, label_h, current_group, size=13,
                          bold=True, color=ACCENT_COLOR)
                y += label_h + Inches(0.15)

        indent = Inches(0.15)  # 그룹에 속한 방문지는 소제목 아래 있다는 느낌을 주기 위해 살짝 들여씀
        placed_any = False
        while idx < len(destinations):
            dest = destinations[idx]
            group = dest.get("region_tag") or None
            is_new_group = bool(group) and group != current_group
            title_w = CONTENT_W - (indent if group else Inches(0))
            title_h = estimate_text_height(dest.get("title", ""), 15, title_w, bold=True)
            image_h = Inches(0.45)
            desc_h = estimate_text_height(dest.get("description", ""), 12, title_w)
            group_label_h = Inches(0)
            if is_new_group:
                group_label_h = estimate_text_height(group, 13, CONTENT_W, bold=True) + Inches(0.15)
            block_h = group_label_h + title_h + Inches(0.06) + image_h + Inches(0.1) \
                + desc_h + Inches(0.18)

            if placed_any and y + block_h > flow.bottom_limit:
                break  # 이 슬라이드엔 더 안 들어감 -> 다음 슬라이드로

            if is_new_group:
                label_h = group_label_h - Inches(0.15)
                add_text(slide, MARGIN, y, CONTENT_W, label_h, group, size=13,
                          bold=True, color=ACCENT_COLOR)
                y += group_label_h
                current_group = group

            x = MARGIN + (indent if group else Inches(0))
            add_text(slide, x, y, title_w, title_h, dest.get("title", ""),
                      size=15, bold=True)
            y += title_h + Inches(0.06)
            add_small_image_placeholder(slide, y, Inches(1.6), image_h, "이미지")
            y += image_h + Inches(0.1)
            add_text(slide, x, y, title_w, desc_h, dest.get("description", ""),
                      size=12, color=MUTED_COLOR)
            y += desc_h + Inches(0.18)

            placed_any = True
            idx += 1

        flow.y = y
        slides.append(slide)
    return slides


def build_background_slide(flow, background_story):
    """배경 이야기 섹션. 예전엔 "OOO란?" 형태의 kicker 소제목을 title 위에 따로
    붙였는데, 모든 상품 기획안에 기계적으로 반복되는 상투적 표현이라 제거함
    (title만으로 바로 임팩트 있게 시작 — prompt_builder.py도 함께 수정됨).

    title은 20pt 굵은 글씨라 길면 2줄로 줄바꿈되는데, 예전엔 고정 Inches(0.55)만
    확보해두고 그 아래 content를 바로 이어 그려서, title이 2줄이 되는 순간 content
    첫 줄과 겹치는 버그가 있었다(카라코람 가을 버전 테스트에서 확인됨 — "실크로드의
    마지막 관문, 카라코람이 품은 가을의 황금빛"이 2줄로 줄바꿈되며 바로 아래 문단과
    겹쳤음). build_experience_slide/build_safety_slide와 같은 문제라 같은 방식
    (estimate_text_height로 실제 높이 추정)으로 고친다."""
    if not background_story:
        return None
    title = background_story.get("title", "")
    content = background_story.get("content", "")
    title_h = estimate_text_height(title, 20, CONTENT_W, bold=True) if title else Inches(0)
    content_h = estimate_text_height(content, 12, CONTENT_W)
    total_h = (title_h + Inches(0.1) if title else Inches(0)) + content_h

    y = flow.ensure(total_h)
    slide = flow.slide
    if title:
        add_text(slide, MARGIN, y, CONTENT_W, title_h, title,
                  size=20, bold=True, align=PP_ALIGN.CENTER)
        y += title_h + Inches(0.1)
    add_text(slide, MARGIN, y, CONTENT_W, content_h, content, size=12, align=PP_ALIGN.CENTER)
    y += content_h
    flow.y = y
    return slide


def build_reasons_slide(flow, why_reasons, product_name=""):
    """'왜 사천성인가' 같은 이유 N가지 섹션.
    타이틀은 AI에게 맡기지 않고 "{상품명} 포인트 0N" 형태로 코드에서 자동 생성한다
    (개수 기반 기계적 표기라 AI보다 코드가 더 정확함)."""
    if not why_reasons:
        return None
    heading = f"{product_name} 포인트 {len(why_reasons):02d}".strip()
    reason_blocks = []
    total_h = Inches(0.6)
    for reason in why_reasons:
        title_h = estimate_text_height(reason.get("title", ""), 15, CONTENT_W, bold=True)
        content_h = estimate_text_height(reason.get("content", ""), 12, CONTENT_W)
        reason_blocks.append((title_h, content_h))
        total_h += title_h + Inches(0.07) + content_h + Inches(0.22)

    y = flow.ensure(total_h)
    slide = flow.slide
    add_section_bar(slide, y, heading)
    y += Inches(0.6)
    for reason, (title_h, content_h) in zip(why_reasons, reason_blocks):
        add_text(slide, MARGIN, y, CONTENT_W, title_h, reason.get("title", ""),
                  size=15, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
        y += title_h + Inches(0.07)
        add_text(slide, MARGIN, y, CONTENT_W, content_h, reason.get("content", ""),
                  size=12, align=PP_ALIGN.CENTER)
        y += content_h + Inches(0.22)
    flow.y = y
    return slide


def build_transport_slide(flow, transport_spec):
    """열차/크루즈처럼 이동수단 자체가 상품의 핵심 매력인 경우의 스펙 섹션.
    안나푸르나(2296 남극 크루즈), 호주 더 간 열차(1827) 상품설명 이미지 분석에서
    반복 확인된 "이동수단 스펙표"(객실타입/부대시설/톤수/안전등급 등) 패턴을 반영."""
    if not transport_spec or not transport_spec.get("specs"):
        return None
    image_h = Inches(0.95)
    rows = []
    total_h = (Inches(0.55) if transport_spec.get("title") else Inches(0)) + image_h
    for spec in transport_spec["specs"]:
        label = spec.get("label", "")
        value = spec.get("value", "")
        row_h = estimate_text_height(f"{label}: {value}", 12, CONTENT_W)
        rows.append((label, value, row_h))
        total_h += row_h + Inches(0.07)

    y = flow.ensure(total_h)
    slide = flow.slide
    if transport_spec.get("title"):
        add_section_bar(slide, y, transport_spec["title"])
        y += Inches(0.55)
    add_small_image_placeholder(slide, y, Inches(2.2), Inches(0.85), "이동수단 이미지")
    y += image_h
    for label, value, row_h in rows:
        add_text(slide, MARGIN, y, Inches(1.6), row_h, label, size=12, bold=True, color=ACCENT_COLOR)
        add_text(slide, MARGIN + Inches(1.7), y, CONTENT_W - Inches(1.7), row_h, value, size=12)
        y += row_h + Inches(0.07)
    flow.y = y
    return slide


def build_guide_slide(flow, guide_profile):
    """인솔자/가이드/담당 임원 프로필 섹션. 제주도 가이드 이력, 산티아고 인솔자
    경력 카드, 트레킹(킬리만자로 40회 등정 임원) 상품설명 이미지에서 반복 확인된
    "회사 구성원 신뢰 요소" 패턴을 반영."""
    if not guide_profile:
        return None
    header_h = Inches(0.55)
    bio_heights = []
    total_h = header_h
    for guide in guide_profile:
        bio_h = estimate_text_height(guide.get("bio", ""), 11, CONTENT_W)
        bio_heights.append(bio_h)
        total_h += Inches(1.18) + Inches(0.25) + bio_h + Inches(0.18)

    y = flow.ensure(total_h)
    slide = flow.slide
    add_section_bar(slide, y, "함께하는 사람들")
    y += header_h
    for guide, bio_h in zip(guide_profile, bio_heights):
        add_small_image_placeholder(slide, y, Inches(1.1), Inches(1.1), "프로필 사진")
        y += Inches(1.18)
        name_title = guide.get("name", "")
        if guide.get("title"):
            name_title = f"{name_title}  ({guide['title']})" if name_title else guide["title"]
        add_text(slide, MARGIN, y, CONTENT_W, Inches(0.28), name_title, size=13, bold=True,
                  align=PP_ALIGN.CENTER)
        y += Inches(0.25)
        add_text(slide, MARGIN, y, CONTENT_W, bio_h, guide.get("bio", ""), size=11,
                  color=MUTED_COLOR, align=PP_ALIGN.CENTER)
        y += bio_h + Inches(0.18)
    flow.y = y
    return slide


def build_meal_slide(flow, meal_info):
    """"트레킹/여행 중 식사는 어떻게 하나요?" 실용 정보 Q&A 섹션. 일본알프스,
    키르기즈스탄, 마칼루, 하얼빈 등 지역이 전혀 다른 다수 상품에서 반복 확인된
    패턴으로, safety_note와 동일한 question/answer 구조를 재사용."""
    if not meal_info or not meal_info.get("question"):
        return None
    question_h = estimate_text_height(meal_info["question"], 15, CONTENT_W, bold=True)
    ans_h = estimate_text_height(meal_info.get("answer", ""), 12, CONTENT_W)
    total_h = question_h + Inches(0.1) + ans_h

    y = flow.ensure(total_h)
    slide = flow.slide
    add_text(slide, MARGIN, y, CONTENT_W, question_h, meal_info["question"], size=15,
              bold=True, align=PP_ALIGN.CENTER)
    y += question_h + Inches(0.1)
    add_text(slide, MARGIN, y, CONTENT_W, ans_h, meal_info.get("answer", ""), size=12,
              align=PP_ALIGN.CENTER)
    y += ans_h
    flow.y = y
    return slide


def build_route_compare_slide(flow, route_compare):
    """두 노선/코스를 비교하는 표 섹션"""
    if not route_compare or not route_compare.get("routes"):
        return None
    routes = route_compare["routes"]
    row_h = Inches(0.9)
    header_h = Inches(0.6) if route_compare.get("title") else Inches(0)
    total_h = header_h + Inches(0.5) + row_h * 4  # 이름줄(0.5) + 기준 4행

    y = flow.ensure(total_h)
    slide = flow.slide
    if route_compare.get("title"):
        add_section_bar(slide, y, route_compare["title"])
        y += Inches(0.6)
    col_w = CONTENT_W / len(routes)
    criteria = ["course", "scenery", "appeal", "summary"]
    criteria_label = {"course": "코스", "scenery": "풍경", "appeal": "매력", "summary": "한줄 요약"}
    for ri, route in enumerate(routes):
        x = MARGIN + col_w * ri
        add_text(slide, x, y, col_w, Inches(0.4), route.get("name", ""), size=14, bold=True,
                  align=PP_ALIGN.CENTER, color=ACCENT_COLOR)
    y += Inches(0.5)
    for crit in criteria:
        for ri, route in enumerate(routes):
            x = MARGIN + col_w * ri
            val = route.get(crit, "")
            add_text(slide, x, y, col_w, row_h, f"[{criteria_label[crit]}]\n{val}", size=10,
                      align=PP_ALIGN.CENTER)
        y += row_h
    flow.y = y
    return slide


def build_experience_slide(flow, brand_tagline, experience_points):
    """브랜드 소구 문구 + 경험 포인트(아이콘 카드 N개).
    예전엔 brand_points(불릿 목록)를 따로 받아 여기 같이 나열했는데, experience_points와
    내용이 거의 그대로 중복되는 문제가 있어(예: '노쇼핑/노옵션'이 두 번 나옴) brand_points는
    제거하고 experience_points 하나로 통일한다."""
    if not brand_tagline and not experience_points:
        return None
    tagline_h = estimate_text_height(brand_tagline, 16, CONTENT_W, bold=True) if brand_tagline else Inches(0)
    col_w = CONTENT_W / len(experience_points) if experience_points else CONTENT_W
    title_h = Inches(0)
    desc_h = Inches(0)
    if experience_points:
        title_h = max(
            estimate_text_height(ep.get("title", ""), 12, col_w - Inches(0.1), bold=True)
            for ep in experience_points
        )
        desc_h = max(
            estimate_text_height(ep.get("description", ""), 10, col_w - Inches(0.1))
            for ep in experience_points
        )
    # 카드 제목이 1줄일 때만 맞는 고정 간격(예전 Inches(0.3))을 쓰면, 제목이 2줄로
    # 줄바꿈되는 순간 바로 아래 설명 텍스트와 겹치는 버그가 있었음 — title_h를 실제
    # 추정 높이로 계산해서 다음 요소를 그만큼 아래로 밀어내도록 수정.
    total_h = (tagline_h + Inches(0.2) if brand_tagline else Inches(0)) \
        + (Inches(0.65) + title_h + Inches(0.1) + desc_h if experience_points else Inches(0))

    y = flow.ensure(total_h)
    slide = flow.slide
    if brand_tagline:
        add_text(slide, MARGIN, y, CONTENT_W, tagline_h, brand_tagline, size=16, bold=True,
                  align=PP_ALIGN.CENTER)
        y += tagline_h + Inches(0.2)
    if experience_points:
        for i, ep in enumerate(experience_points):
            x = MARGIN + col_w * i
            add_image_placeholder(slide, x + Inches(0.05), y, col_w - Inches(0.1), Inches(0.55), "아이콘")
        y += Inches(0.65)
        for i, ep in enumerate(experience_points):
            x = MARGIN + col_w * i
            add_text(slide, x, y, col_w - Inches(0.1), title_h, ep.get("title", ""), size=12,
                      bold=True, align=PP_ALIGN.CENTER)
        y += title_h + Inches(0.1)
        for i, ep in enumerate(experience_points):
            x = MARGIN + col_w * i
            dh = estimate_text_height(ep.get("description", ""), 10, col_w - Inches(0.1))
            add_text(slide, x, y, col_w - Inches(0.1), dh, ep.get("description", ""), size=10,
                      color=MUTED_COLOR, align=PP_ALIGN.CENTER)
        y += desc_h
    flow.y = y
    return slide


def build_highlights_slides(flow, highlights, heading=None):
    """번호 매긴 여정 하이라이트 카드 (destinations와 별개 — 더 큰 테마 단위)"""
    if not highlights:
        return []
    heading = heading or "여정 하이라이트"  # AI가 빠뜨려도 타이틀 없는 슬라이드가 나가지 않도록 기본값
    slides = []
    idx = 0
    first = True
    header_h = Inches(0.55)

    while idx < len(highlights):
        if first:
            y = flow.ensure(header_h + Inches(0.5))
            slide = flow.slide
            add_section_bar(slide, y, heading)
            y += header_h
            first = False
        else:
            slide = flow.new_slide()
            y = flow.y

        placed_any = False
        while idx < len(highlights):
            item = highlights[idx]
            num_label = f"{idx + 1:02d}"
            title_h = estimate_text_height(item.get("title", ""), 14, CONTENT_W, bold=True)
            image_h = Inches(0.45)
            desc_h = estimate_text_height(item.get("description", ""), 11, CONTENT_W)
            block_h = Inches(0.22) + title_h + Inches(0.06) + image_h + Inches(0.1) + desc_h + Inches(0.18)
            if placed_any and y + block_h > flow.bottom_limit:
                break
            add_text(slide, MARGIN, y, Inches(0.6), Inches(0.25), num_label, size=13, bold=True,
                      color=ACCENT_COLOR)
            y += Inches(0.26)
            add_text(slide, MARGIN, y, CONTENT_W, title_h, item.get("title", ""), size=14, bold=True)
            y += title_h + Inches(0.06)
            add_small_image_placeholder(slide, y, Inches(1.6), image_h, "이미지")
            y += image_h + Inches(0.1)
            add_text(slide, MARGIN, y, CONTENT_W, desc_h, item.get("description", ""), size=11,
                      color=MUTED_COLOR)
            y += desc_h + Inches(0.18)
            placed_any = True
            idx += 1

        flow.y = y
        slides.append(slide)
    return slides


def build_season_slide(flow, season, season_table=None):
    if not season or (not season.get("content") and not season_table):
        return None
    header_h = Inches(0.6)
    stat_h = Inches(0.5) if season.get("stat_line") else Inches(0)
    content_h = estimate_text_height(season.get("content", ""), 12, CONTENT_W)
    table_h = Inches(0.85) + Inches(0.3) if season_table else Inches(0)
    total_h = header_h + stat_h + content_h + Inches(0.3) + table_h

    y = flow.ensure(total_h)
    slide = flow.slide
    add_section_bar(slide, y, season.get("title", "언제 가면 좋을까?"))
    y += header_h
    if season.get("stat_line"):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, CONTENT_W, Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_COLOR
        bar.line.fill.background()
        _tf_setup(bar.text_frame, season["stat_line"], 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        y += Inches(0.5)
    add_text(slide, MARGIN, y, CONTENT_W, content_h, season.get("content", ""), size=12)
    y += content_h + Inches(0.3)
    if season_table:
        add_image_placeholder(slide, MARGIN, y, CONTENT_W, Inches(0.7),
                               "월별 기온 차트 (자리표시 — 실제 그래픽은 디자이너 작업)")
        y += Inches(0.85)
        header = "  ".join(f"{row.get('month','')}" for row in season_table)
        add_text(slide, MARGIN, y, CONTENT_W, Inches(0.3), header, size=10, color=MUTED_COLOR,
                  align=PP_ALIGN.CENTER)
        y += Inches(0.3)
    flow.y = y
    return slide


def build_safety_slide(flow, altitude_profile, safety_note):
    """경유지 고도 프로필(있는 경우) + 안전/난이도 관련 표준 안내.
    고산 트레킹의 '고산증', 도보순례의 '체력/보험', 일반 하이킹의 '난이도' 등
    카테고리에 따라 톤이 다른 표준 안내문을 담는 범용 섹션."""
    if not altitude_profile and not safety_note:
        return None
    ans_h = Inches(0)
    question_h = Inches(0)
    qa_h = Inches(0)
    if safety_note and safety_note.get("question"):
        question_h = estimate_text_height(safety_note["question"], 15, CONTENT_W, bold=True)
        ans_h = estimate_text_height(safety_note.get("answer", ""), 12, CONTENT_W)
        qa_h = question_h + Inches(0.1) + ans_h + Inches(0.35)

    n = len(altitude_profile) if altitude_profile else 0
    col_w = CONTENT_W / max(n, 1)
    gap = Inches(0.06)
    labels = []
    if altitude_profile:
        for stop in altitude_profile:
            # 자리표시 박스는 "숙박"이 아니라 "지도"여야 함 — 예전 코드가 숙박/롯지
            # 소개 슬라이드에서 라벨만 안 바꾼 채 복붙된 흔적으로 보임.
            label = f"{stop.get('name','')}\n{stop.get('altitude','')}"
            extra = " / ".join(v for v in (stop.get("distance"), stop.get("duration")) if v)
            if extra:
                label += f"\n{extra}"
            if stop.get("highlight"):
                label += f"\n{stop['highlight']}"
            labels.append(label)
    # 라벨 줄 수가 highlight 유무에 따라 2~4줄로 달라지므로, 고정 Inches(0.65) 대신
    # 실제 텍스트 높이를 추정해서 다음 요소와 겹치지 않게 함(build_experience_slide와
    # 동일한 문제였음).
    label_h = max((estimate_text_height(l, 9, col_w - gap * 2) for l in labels), default=Inches(0.5))
    profile_h = (Inches(0.4) + Inches(0.6) + label_h) if altitude_profile else Inches(0)
    total_h = qa_h + profile_h

    y = flow.ensure(total_h)
    slide = flow.slide
    if safety_note and safety_note.get("question"):
        add_text(slide, MARGIN, y, CONTENT_W, question_h, safety_note["question"], size=15,
                  bold=True, align=PP_ALIGN.CENTER)
        y += question_h + Inches(0.1)
        add_text(slide, MARGIN, y, CONTENT_W, ans_h, safety_note.get("answer", ""), size=12,
                  align=PP_ALIGN.CENTER)
        y += ans_h + Inches(0.35)
    if altitude_profile:
        add_text(slide, MARGIN, y, CONTENT_W, Inches(0.3),
                  "구간별 고도 프로필 (자리표시 — 실제 그래픽은 디자이너 작업)",
                  size=10, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
        y += Inches(0.4)
        for i in range(n):
            x = MARGIN + col_w * i
            add_image_placeholder(slide, x + gap, y, col_w - gap * 2, Inches(0.5), "지도")
        y += Inches(0.6)
        for i, label in enumerate(labels):
            x = MARGIN + col_w * i
            add_text(slide, x + gap, y, col_w - gap * 2, label_h, label, size=9, align=PP_ALIGN.CENTER)
        y += label_h
    flow.y = y
    return slide


def build_banner_request_slide(flow, cover, banner_copy=None):
    """배너 기획 페이지 — 실제 회사 배너제작 템플릿(배너제작.pptx)의 레이아웃을
    그대로 재현. 4개 배너 슬롯(메인 와이드/서브메인 띠/서브메인 2단/지역 리스트)에
    각각 이미지 자리와 '태그라인+타이틀' 텍스트를 넣는다. 스펙 라벨/안내 문구는
    회사 표준이라 고정값이며, 지역/상품과 무관하게 항상 그대로 포함. 절대 위치로
    실제 배너 템플릿과 맞춰야 해서 다른 섹션과 공유하지 않고 항상 전용 슬라이드로 만든다.

    예전엔 cover.tagline+product_name(표지 전체 카피)을 4개 슬롯에 그대로 복붙해서
    문구가 길고 밋밋했음 — 실제 회사 배너는 표지 카피와 다르게, 상품에서 가장
    눈에 띄는 포인트만 뽑아 훨씬 짧고 강렬하게 압축한 별도 카피를 쓴다(첨부 배너
    예시: "봄으로 물든 차마고도를 걷다 / 카피할 수 없는 오리지널의 위엄" +
    "세계 3대 트레킹 / 호도협·옥룡설산"). banner_copy가 있으면 그걸 쓰고, AI가
    누락했을 때만 cover 필드로 대체(하위 호환).

    4개 슬롯 모두 "부제(후킹 카피) + 상품명(타이틀)" 두 요소를 함께 넣는다. 처음엔
    상품명을 부제보다 작게 눌러서 부제(후킹 문구)가 시선을 먼저 받게 했었는데,
    실제 배너에서는 반대로 상품명(타이틀)이 메인 카피이고 부제는 그걸 보조하는
    작은 문구라는 피드백을 받아 위계를 뒤집었다 — 부제는 작고 옅게, 상품명은
    크고 굵게 렌더링한다. 부제와 상품명을 하나의 멀티라인 텍스트박스에 같은
    크기로 합쳐 넣으면(예전 방식) 위계가 전혀 안 생기고, 심지어 메인 와이드
    배너는 상품명 자체가 빠지는 버그가 있었다(카라코람 테스트에서 확인됨) —
    부제와 상품명을 별도 텍스트박스로 쌓아 크기를 다르게 주고, 4개 슬롯 전부에
    상품명을 반드시 포함시킨다."""
    slide = flow.new_slide()
    banner_copy = banner_copy or {}
    kicker = banner_copy.get("kicker") or cover.get("tagline", "")
    title = banner_copy.get("title") or cover.get("product_name", "")
    product_name = cover.get("product_name", "")

    def _add_banner_copy(left, top, width, subtitle_text, subtitle_size, name_size):
        """부제(작고 옅은 글씨)를 먼저 쌓고, 그 아래 상품명/타이틀(더 크고 굵은
        글씨)을 이어 쌓는다 — 실제 들어간 줄 수에 맞춰 높이를 추정해서 다음
        요소랑 안 겹치게 한다."""
        y = top
        if subtitle_text:
            h = estimate_text_height(subtitle_text, subtitle_size, width)
            add_text(slide, left, y, width, h, subtitle_text, size=subtitle_size,
                      color=MUTED_COLOR, align=PP_ALIGN.CENTER)
            y += h
        if product_name:
            h = estimate_text_height(product_name, name_size, width, bold=True)
            add_text(slide, left, y, width, h, product_name, size=name_size,
                      bold=True, align=PP_ALIGN.CENTER)
            y += h
        return y

    add_section_bar(slide, Inches(0), "배너 기획", height=Inches(0.47), size=14)
    add_text(slide, Inches(0.106), Inches(0.675), Inches(1.717), Inches(0.404),
              "배너제작요청", size=13, bold=True)
    add_text(slide, Inches(1.823), Inches(0.733), Inches(4.672), Inches(0.303),
              "홈페이지 개편에 따라, 배너 디자인이 전면 교체되었습니다.", size=9, color=MUTED_COLOR)

    # 메인 와이드 배너 — 부제 2줄(kicker+title, 후킹 카피, size 11) + 상품명(size 18, 굵게)
    add_text(slide, Inches(0.106), Inches(1.271), Inches(5.701), Inches(0.303),
              "메인 와이드 배너 (PC: 1920x700, MO: 750x510), 가이드라인에 맞춰 제작", size=9, color=MUTED_COLOR)
    add_image_placeholder(slide, Inches(0.217), Inches(1.630), Inches(5.475), Inches(1.817), "이미지")
    _add_banner_copy(Inches(0.388), Inches(2.158), Inches(5.133),
                      f"{kicker}\n{title}", subtitle_size=11, name_size=18)

    # 서브메인 띠배너 — 부제 1줄(kicker, size 11) + 상품명(size 16, 굵게)
    add_text(slide, Inches(0.081), Inches(3.581), Inches(5.642), Inches(0.303),
              "서브메인 띠배너 (PC: 1920x200, MO: 750x200), 가이드라인에 맞춰 제작", size=9, color=MUTED_COLOR)
    add_image_placeholder(slide, Inches(0.136), Inches(4.021), Inches(7.028), Inches(0.979), "이미지")
    _add_banner_copy(Inches(1.184), Inches(4.173), Inches(5.133),
                      kicker, subtitle_size=11, name_size=16)

    # 서브메인 2단배너 — 부제 1줄(kicker, size 9) + 상품명(size 13, 굵게)
    add_text(slide, Inches(0.136), Inches(5.137), Inches(5.642), Inches(0.303),
              "서브메인 2단배너 (PC: 590x370, MO: 585x670), 가이드라인에 맞춰 제작", size=9, color=MUTED_COLOR)
    add_image_placeholder(slide, Inches(0.221), Inches(5.521), Inches(2.835), Inches(1.771), "이미지")
    _add_banner_copy(Inches(0.288), Inches(6.138), Inches(2.835),
                      kicker, subtitle_size=9, name_size=13)

    # 지역 리스트 배너 — 부제 1줄(kicker, size 11) + 상품명(size 16, 굵게)
    add_text(slide, Inches(0.205), Inches(7.481), Inches(5.701), Inches(0.303),
              "지역 리스트 배너 (PC: 1200x207, MO: 750x207), 가이드라인에 맞춰 제작", size=9, color=MUTED_COLOR)
    add_image_placeholder(slide, Inches(0.316), Inches(7.989), Inches(6.871), Inches(1.336), "이미지")
    _add_banner_copy(Inches(1.184), Inches(8.280), Inches(5.133),
                      kicker, subtitle_size=11, name_size=16)

    return slide


def build_review_notes_slide(flow, review):
    """Gemini 교차검수 결과를 PPTX 마지막 페이지에 그대로 남긴다 — 예전엔 Streamlit
    화면에만 표시되고 다운로드한 PPTX에는 안 남아서, 파일만 디자이너/다른 담당자에게
    넘기면 검수 내역이 통째로 사라지는 문제가 있었다. review가 없거나 검수 자체를
    건너뛴 경우(_dry_run — GEMINI_API_KEY 미설정)는 내용 없는 안내 페이지를 억지로
    넣지 않기 위해 슬라이드를 만들지 않는다."""
    if not review or review.get("_dry_run"):
        return []
    issues = review.get("issues") or []

    slide = flow.new_slide()
    y = flow.y
    add_section_bar(slide, y, "Gemini 교차검수 결과 (내부 참고용)")
    y += Inches(0.55)
    slides = [slide]

    if not issues:
        add_text(slide, MARGIN, y, CONTENT_W, Inches(0.4),
                  "✅ 검수 통과 — 왜곡/날조·저작권/표절·사실확인·비문/맞춤법 이슈가 "
                  "발견되지 않았습니다.", size=12, align=PP_ALIGN.CENTER)
        flow.y = y + Inches(0.4)
        return slides

    if review.get("summary"):
        summary_h = estimate_text_height(review["summary"], 11, CONTENT_W)
        add_text(slide, MARGIN, y, CONTENT_W, summary_h, review["summary"], size=11,
                  color=MUTED_COLOR)
        y += summary_h + Inches(0.2)

    for issue in issues:
        header_text = f"[{issue.get('category', '')} · {issue.get('severity', '')}] {issue.get('field', '')}"
        header_h = estimate_text_height(header_text, 12, CONTENT_W, bold=True)
        quote = issue.get("quote", "")
        quote_h = estimate_text_height(quote, 10, CONTENT_W - Inches(0.2)) if quote else Inches(0)
        explanation_h = estimate_text_height(issue.get("explanation", ""), 11, CONTENT_W)
        block_h = header_h + Inches(0.05) \
            + (quote_h + Inches(0.05) if quote else Inches(0)) \
            + explanation_h + Inches(0.22)

        if y + block_h > flow.bottom_limit:
            slide = flow.new_slide()
            y = flow.y
            slides.append(slide)

        add_text(slide, MARGIN, y, CONTENT_W, header_h, header_text, size=12, bold=True,
                  color=ACCENT_COLOR)
        y += header_h + Inches(0.05)
        if quote:
            add_text(slide, MARGIN + Inches(0.1), y, CONTENT_W - Inches(0.2), quote_h,
                      f"“{quote}”", size=10, color=MUTED_COLOR)
            y += quote_h + Inches(0.05)
        add_text(slide, MARGIN, y, CONTENT_W, explanation_h, issue.get("explanation", ""), size=11)
        y += explanation_h + Inches(0.22)

    flow.y = y
    return slides


def build(content_json, out_path, review=None):
    """content_json(정현지 스키마) -> 새 PPTX 파일 생성.
    review(reviewer.review_content()의 반환값)를 넘기면 마지막에 검수 결과 페이지를
    덧붙인다 — 넘기지 않으면(기본값 None) 예전과 동일하게 검수 페이지 없이 생성된다."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    flow = SlideFlow(prs)

    cover = content_json.get("cover", {})
    build_cover_slide(flow, cover, content_json.get("watermark_label", ""))
    build_background_slide(flow, content_json.get("background_story"))
    build_reasons_slide(flow, content_json.get("why_reasons"), product_name=cover.get("product_name", ""))
    build_destination_slides(
        flow,
        content_json.get("destinations", []),
        section_title=content_json.get("destinations_heading"),
        theme_line=None,
    )
    # tour_spots: 트레킹 상품 중 트레킹 코스 + 관광(도시/유적/박물관) 코스가 함께 있는
    # 경우에만 채워진다 — 카라코람 테스트에서 훈자 마을/이슬라마바드/탁실라 박물관 같은
    # "관광" 코스가 트레킹 코스(destinations)와 구분 없이 한 섹션에 섞여 나온 문제를
    # 고치기 위해 별도 섹션으로 분리했다. build_destination_slides와 카드 레이아웃은
    # 동일하고 헤딩만 다르다.
    build_destination_slides(
        flow,
        content_json.get("tour_spots", []),
        section_title=content_json.get("tour_spots_heading"),
        theme_line=None,
    )
    build_route_compare_slide(flow, content_json.get("route_compare"))
    build_transport_slide(flow, content_json.get("transport_spec"))
    build_experience_slide(
        flow,
        content_json.get("brand_tagline", ""),
        content_json.get("experience_points"),
    )
    build_guide_slide(flow, content_json.get("guide_profile"))
    # highlights/highlights_heading 필드는 제거됨 — why_reasons("포인트 0N" 섹션)와
    # 지시문이 실질적으로 같은 "상품 차별점/테마 하이라이트" 내용을 요구해, AI가 같은
    # 내용을 문구만 바꿔 두 번 반복하는 문제가 반복 발생했다(예: 멕시코 문명기행
    # 테스트에서 "칸쿤 없이 완성하는 내륙 문명 루트" 등 4개 항목이 포인트 섹션과
    # 하이라이트 섹션에 그대로 중복). 스키마에서 highlights를 없애 why_reasons
    # 하나로 통일한다(build_highlights_slides 함수 자체는 향후 필요시를 위해 남겨둠).
    build_season_slide(flow, content_json.get("season", {}), content_json.get("season_table"))
    build_meal_slide(flow, content_json.get("meal_info"))
    build_safety_slide(flow, content_json.get("altitude_profile"), content_json.get("safety_note"))
    build_banner_request_slide(flow, cover, content_json.get("banner_copy"))
    build_review_notes_slide(flow, review)
    prs.save(out_path)
    return prs
