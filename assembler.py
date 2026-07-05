"""검증된 텍스트 주입 로직 — JSON 필드를 template_map 기준으로 PPTX 도형에 씀"""
import json
from copy import deepcopy
from pptx import Presentation

NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

def _set_text(slide, shape_id, new_text):
    for shape in slide.shapes:
        if shape.shape_id != shape_id or not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            return "SKIPPED(no base run)"
        lines = str(new_text).split("\n")
        base_para_xml = deepcopy(tf.paragraphs[0]._p)
        txBody = tf._txBody
        for p in list(tf.paragraphs):
            txBody.remove(p._p)
        for line in lines:
            new_p = deepcopy(base_para_xml)
            r_elems = new_p.findall(f'.//{NS}r')
            if r_elems:
                keep = r_elems[0]
                for extra in r_elems[1:]:
                    extra.getparent().remove(extra)
                t_elem = keep.find(f'{NS}t')
                t_elem.text = line
            txBody.append(new_p)
        return "OK"
    return "NOT FOUND"


def _delete_shape(slide, shape_id):
    """슬롯이 남을 때(실제 데이터가 템플릿 슬롯 수보다 적을 때) 텍스트만 지우지 않고
    도형 자체를 슬라이드에서 제거한다. 텍스트만 비우면 빈 상자/잔여 이미지가 남아
    레이아웃이 깨지므로, pptx 스킬 가이드대로 그룹 전체를 삭제한다."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            shape._element.getparent().remove(shape._element)
            return True
    return False

def _get_nested(data, path):
    """'why_hyecho.points.0' 같은 경로 문자열로 JSON에서 값 추출"""
    cur = data
    for key in path.split("."):
        if isinstance(cur, list):
            idx = int(key)
            if idx >= len(cur):
                return None
            cur = cur[idx]
        elif isinstance(cur, dict):
            cur = cur.get(key)
        else:
            return None
        if cur is None:
            return None
    return cur

def _scan_for_leaks(prs, leak_terms):
    """조립이 끝난 뒤, 템플릿 원본(예: 안데스/방글라데시)에 있던 고유명사가
    아직 남아있는 도형이 있는지 전수 스캔한다. 지금까지 매핑하지 못한 도형이든,
    앞으로 완전히 새로운 지역 자료가 들어와서 생기는 예상 못한 누수든,
    조용히 잘못된 결과물이 나가는 대신 여기서 반드시 걸리게 하기 위한 안전장치."""
    if not leak_terms:
        return []
    leaks = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text.strip():
                continue
            for term in leak_terms:
                if term in text:
                    leaks.append((si, shape.shape_id, term, text.strip()[:50]))
                    break
    return leaks


def assemble(content_json, writer_style, template_map_path, out_path):
    with open(template_map_path, encoding="utf-8") as f:
        tmap = json.load(f)[writer_style]

    prs = Presentation(tmap["template"])
    log = []

    # 1) 단일 필드 매핑 (표지, why_hyecho 등 — 기존 로직)
    for field_path, (slide_idx, shape_id) in tmap.get("field_map", {}).items():
        value = _get_nested(content_json, field_path)
        if value is None:
            log.append((field_path, "NO DATA"))
            continue
        # points/bullets 등 리스트나 dict면 한 줄로 합침
        if isinstance(value, dict):
            value = "\n".join(str(v) for v in value.values())
        elif isinstance(value, list):
            value = "\n".join(str(v) for v in value)
        result = _set_text(prs.slides[slide_idx], shape_id, value)
        log.append((field_path, result))

    # 2) 반복 슬롯 그룹 (목적지 목록, 경유지 라벨 등 — 실제 개수가
    #    템플릿 슬롯 수보다 적으면 남는 슬롯의 도형을 통째로 삭제한다)
    for group_name, group_cfg in tmap.get("repeatable_groups", {}).items():
        items = content_json.get(group_name, []) or []
        slots = group_cfg.get("slots", [])

        if len(items) > len(slots):
            log.append((f"{group_name}(개수 초과)",
                        f"WARNING: 실제 {len(items)}개인데 템플릿 슬롯은 {len(slots)}개뿐 — "
                        f"뒤 {len(items) - len(slots)}개는 버려짐. 템플릿 확장 필요"))

        for i, slot in enumerate(slots):
            is_scalar_slot = "shape" in slot  # map_stops처럼 원소가 dict가 아니라 그냥 문자열인 슬롯
            all_shape_refs = dict(slot.get("fields", {}))
            if is_scalar_slot:
                all_shape_refs["value"] = slot["shape"]
            all_shape_refs.update({f"__extra_{j}": ref for j, ref in enumerate(slot.get("extra_shapes", []))})

            if i < len(items):
                item = items[i]
                if is_scalar_slot:
                    slide_idx, shape_id = slot["shape"]
                    label = f"{group_name}[{i}]"
                    if item is None:
                        log.append((label, "NO DATA"))
                    else:
                        result = _set_text(prs.slides[slide_idx], shape_id, item)
                        log.append((label, result))
                else:
                    for field_name, (slide_idx, shape_id) in slot.get("fields", {}).items():
                        value = item.get(field_name) if isinstance(item, dict) else None
                        label = f"{group_name}[{i}].{field_name}"
                        if value is None:
                            log.append((label, "NO DATA"))
                            continue
                        result = _set_text(prs.slides[slide_idx], shape_id, value)
                        log.append((label, result))
                # extra_shapes(이미지 자리 등 텍스트가 없는 도형)는 데이터가 있으면 그대로 둔다
            else:
                # 남는 슬롯 — 텍스트 도형 + 부가 도형(이미지 placeholder 등) 전체 삭제
                for field_name, (slide_idx, shape_id) in all_shape_refs.items():
                    label = f"{group_name}[{i}].{field_name}(초과 슬롯)"
                    deleted = _delete_shape(prs.slides[slide_idx], shape_id)
                    log.append((label, "DELETED" if deleted else "DELETE_FAILED(shape not found)"))

    prs.save(out_path)

    # 조립 후 자동 누수 검사 — 템플릿 원본 지명이 남아있으면 로그에 크게 경고로 남긴다
    leak_terms = tmap.get("leak_check_terms", [])
    leaks = _scan_for_leaks(prs, leak_terms)
    for slide_idx, shape_id, term, preview in leaks:
        log.append((
            f"⚠️ LEAK CHECK — slide{slide_idx}/shape{shape_id}",
            f"템플릿 원본 단어 '{term}' 잔존: \"{preview}...\" — 이 도형은 아직 매핑되지 않았습니다"
        ))
    if not leaks and leak_terms:
        log.append(("leak_check", f"OK — 템플릿 원본 고유명사 {len(leak_terms)}개 잔존 없음"))

    return log
