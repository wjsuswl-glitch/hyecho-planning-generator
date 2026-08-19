"""사업부 원본자료(docx/pptx/이미지) 파싱 모듈 — 2단계 계층 구조 지원"""
import re
import io
import base64
import mimetypes
from PIL import Image, ImageOps
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.ns import qn
TOP_MARKER_RE = re.compile(r"^(\*|●|chapter\s*\d+|[0-9]{2}\s)", re.IGNORECASE)
SUB_NUMBER_RE = re.compile(r"^[0-9]\.\s")
def iter_block_items(doc):
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn('w:p'):
            yield Paragraph(child, doc)
        elif child.tag == qn('w:tbl'):
            yield Table(child, doc)
def parse_docx(path):
    doc = Document(path)
    items = list(iter_block_items(doc))
    has_strict_top = any(
        TOP_MARKER_RE.match(it.text.strip()) and len(it.text.strip()) < 40
        for it in items if isinstance(it, Paragraph) and it.text.strip()
    )
    use_subnumber_as_top = not has_strict_top
    sections = {}
    current_key = "header"
    sections[current_key] = []
    current_sub = None
    def add_content(content):
        nonlocal current_sub
        if current_sub is not None:
            current_sub["items"].append(content)
        else:
            sections[current_key].append(content)
    for it in items:
        if isinstance(it, Paragraph):
            t = it.text.strip()
            if not t:
                continue
            style_name = it.style.name if it.style else ""
            is_heading2 = style_name.startswith("Heading 2")
            if TOP_MARKER_RE.match(t) and len(t) < 40:
                current_key, current_sub = t, None
                sections[current_key] = []
            elif SUB_NUMBER_RE.match(t) and len(t) < 40 and use_subnumber_as_top:
                current_key, current_sub = t, None
                sections[current_key] = []
            elif (SUB_NUMBER_RE.match(t) and len(t) < 40) or (is_heading2 and len(t) < 60):
                current_sub = {"subheading": t, "items": []}
                sections[current_key].append(current_sub)
            else:
                add_content(t)
        else:
            rows = [[c.text.strip().replace("\n", " ") for c in row.cells] for row in it.rows]
            add_content({"table": rows})
    return sections
def detect_format_and_draft_copy(sections):
    """유형 판별 + draft_copy(카피 초안) 존재 여부 감지"""
    keys = list(sections.keys())
    has_numbered = any(re.match(r"^[0-9]{2}\s", k) for k in keys)
    has_chapter = any(re.match(r"^chapter", k, re.IGNORECASE) for k in keys)
    has_star = any(k.startswith("*") for k in keys)
    has_dot = any(k.startswith("●") for k in keys)
    if has_numbered or has_chapter:
        fmt = "B"
    elif has_dot:
        fmt = "C"
    elif has_star:
        fmt = "A"
    elif len(keys) <= 1:
        fmt = "D"
    else:
        fmt = "E"
    draft_copy = None
    if fmt == "B":
        for k, v in sections.items():
            if "디자인팀" in k or "카피라이팅" in k or "가이드" in k:
                for block in v:
                    if isinstance(block, dict) and "table" in block:
                        for row in block["table"]:
                            joined = " ".join(row)
                            if "카피" in joined or len(joined) > 20:
                                draft_copy = joined
                                break
    return {"format_type": fmt, "draft_copy": draft_copy}


def parse_pptx(path):
    """PPTX 자료(예: 사업부에서 준 예전 상품 소개 PPT)를 슬라이드별 텍스트/표로 추출.
    docx의 '*'/'●'/번호 마커 같은 계층 구조는 PPT엔 없으므로, 슬라이드 번호를 그대로
    키로 써서 sections와 같은 형태(dict)로 반환한다 — prompt_builder는 dict 구조면
    무엇이든 그대로 JSON화해서 프롬프트에 넣으므로 별도 변환이 필요 없다."""
    from pptx import Presentation
    prs = Presentation(path)
    result = {}
    for i, slide in enumerate(prs.slides, 1):
        items = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                items.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                rows = [[c.text.strip().replace("\n", " ") for c in row.cells]
                        for row in shape.table.rows]
                items.append({"table": rows})
        if items:
            result[f"슬라이드 {i}"] = items
    return result


def encode_image_block(path, max_dimension=2000, jpeg_quality=85):
    """이미지 파일을 Claude API 멀티모달 메시지에 넣을 수 있는 형태로 base64 인코딩.
    generator.py에서 텍스트 프롬프트와 함께 content 리스트에 섞어 보낸다.

    예전엔 원본 파일을 그대로 base64만 인코딩해서 보냈는데, 사업부에서 받은 옛
    기획안 캡처 이미지 중 가로/세로 한 변이 8000px을 넘는 경우가 있어 Claude API가
    "image dimensions exceed max allowed size: 8000 pixels" 오류로 요청 자체를
    거부하는 문제가 있었다(상품 병합 케이스처럼 기존 상품소개 이미지를 원본
    해상도 그대로 올리는 경우 특히 발생하기 쉬움). Pillow로 열어서 긴 변이
    max_dimension(기본 2000px)을 넘으면 비율을 유지한 채 줄이고, 항상 JPEG로
    다시 인코딩해서 media_type 불일치나 팔레트/투명 채널(RGBA, PNG 등) 문제도
    함께 없앤다. 2000px면 API 하드 제한(8000px)에 여유가 크고, 이미지 속 텍스트
    (기존 상품 설명 등)를 읽는 데도 지장이 없는 해상도다. 휴대폰으로 찍은 사진의
    EXIF 방향 정보도 여기서 반영해 회전 문제를 방지한다."""
    img = Image.open(path)
    img = ImageOps.exif_transpose(img)  # 카메라 회전(EXIF Orientation) 그대로 반영
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")  # RGBA/팔레트 모드는 JPEG로 저장 불가

    width, height = img.size
    longest_side = max(width, height)
    if longest_side > max_dimension:
        scale = max_dimension / longest_side
        new_size = (max(1, round(width * scale)), max(1, round(height * scale)))
        img = img.resize(new_size, Image.LANCZOS)

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=jpeg_quality)
    data = base64.b64encode(buf.getvalue()).decode("utf-8")
    return {
        "type": "image",
        "source": {"type": "base64", "media_type": "image/jpeg", "data": data},
    }
